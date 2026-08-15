from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = Path.home() / "Downloads" / "AI_Web_Test_Automation_Viva.pptx"

INK = RGBColor(17, 24, 39)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(15, 118, 110)
AMBER = RGBColor(180, 83, 9)
RED = RGBColor(185, 28, 28)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PAPER


def add_title(slide, kicker, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.34), Inches(9.2), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.text = kicker.upper()
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = TEAL

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.66), Inches(11.2), Inches(0.72))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK

    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.58), Inches(1.36), Inches(10.8), Inches(0.36))
        p = box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED


def add_footer(slide, number):
    box = slide.shapes.add_textbox(Inches(11.55), Inches(6.95), Inches(0.55), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.text = f"{number:02d}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def add_card(slide, x, y, w, h, title, body, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(226, 232, 240)
    shape.adjustments[0] = 0.08

    tag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()

    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.16), Inches(w - 0.38), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = INK

    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.52), Inches(w - 0.38), Inches(h - 0.64))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    p.font.size = Pt(10.5)
    p.font.color.rgb = MUTED


def add_bullets(slide, x, y, items, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(10.7), Inches(4.7))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def add_flow(slide, labels):
    x = 0.72
    y = 2.35
    w = 1.38
    h = 0.72
    for idx, label in enumerate(labels):
        color = [BLUE, TEAL, AMBER, BLUE, TEAL][idx % 5]
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.adjustments[0] = 0.12
        p = shape.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        if idx < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + w + 0.12), Inches(y + 0.17), Inches(0.55), Inches(0.34))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(148, 163, 184)
            arrow.line.fill.background()
        x += w + 0.78


def build():
    prs = Presentation()
    prs.slide_width = Inches(12.8)
    prs.slide_height = Inches(7.2)
    blank = prs.slide_layouts[6]

    slides = []
    for _ in range(10):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        slides.append(slide)

    s = slides[0]
    add_title(s, "Project Viva", "AI-Powered Intelligent Web Test Automation and Analytics Platform", "Dynamic test generation, automated browser execution, evidence capture, and failure analysis")
    add_card(s, 0.72, 2.18, 3.4, 1.45, "Student Focus", "A simple platform that turns natural language testing needs into repeatable browser test runs.", TEAL)
    add_card(s, 4.7, 2.18, 3.4, 1.45, "Core Tools", "FastAPI, Playwright, SQLite, SQLAlchemy, optional Ollama/Qwen AI model.", BLUE)
    add_card(s, 8.68, 2.18, 3.4, 1.45, "Outcome", "Generated test cases, execution evidence, run history, and analytics dashboard.", AMBER)
    add_footer(s, 1)

    s = slides[1]
    add_title(s, "Problem", "Manual web test automation is still slow and difficult.", "Traditional tools require script writing, locator maintenance, log review, and separate reporting.")
    add_bullets(s, 0.95, 2.05, [
        "Frequent web application changes demand continuous testing.",
        "Non-programmers struggle to create and maintain automation scripts.",
        "Execution evidence and failure analysis are often scattered across tools.",
        "The project combines generation, execution, evidence, and analytics in one place.",
    ])
    add_footer(s, 2)

    s = slides[2]
    add_title(s, "Objectives", "The platform converts simple test requirements into browser evidence.", None)
    add_card(s, 0.72, 1.7, 3.45, 1.25, "Generate", "Create structured test scenarios from website URL and natural language requirement.", BLUE)
    add_card(s, 4.68, 1.7, 3.45, 1.25, "Execute", "Run generated steps through Playwright in Chromium and capture status.", TEAL)
    add_card(s, 8.64, 1.7, 3.45, 1.25, "Analyze", "Store results, logs, screenshots, duration, and pass/fail trends.", AMBER)
    add_card(s, 2.7, 3.45, 3.45, 1.25, "Explain", "Summarize failures in readable form for easier debugging.", RED)
    add_card(s, 6.66, 3.45, 3.45, 1.25, "Simplify", "Keep setup and demo flow suitable for academic evaluation.", BLUE)
    add_footer(s, 3)

    s = slides[3]
    add_title(s, "Architecture", "A modular FastAPI system connects AI planning, browser automation, and analytics.", None)
    add_flow(s, ["Dashboard", "FastAPI", "AI Planner", "Playwright", "SQLite"])
    add_card(s, 0.95, 4.05, 3.1, 1.15, "Frontend", "Single-page dashboard for test generation, run control, logs, and metrics.", BLUE)
    add_card(s, 4.85, 4.05, 3.1, 1.15, "Backend", "REST APIs, validation, database persistence, and service orchestration.", TEAL)
    add_card(s, 8.75, 4.05, 3.1, 1.15, "Automation", "Chromium execution with screenshots, console messages, and error capture.", AMBER)
    add_footer(s, 4)

    s = slides[4]
    add_title(s, "Workflow", "Every run follows a clear generation-to-analysis pipeline.", None)
    add_bullets(s, 0.9, 1.75, [
        "1. User enters website URL and testing requirement.",
        "2. AI or fallback planner creates structured steps.",
        "3. Test case is saved for repeat execution.",
        "4. Playwright opens the page and executes each step.",
        "5. Screenshots, logs, status, and duration are stored.",
        "6. Dashboard updates analytics and run history.",
    ], size=13.5)
    add_footer(s, 5)

    s = slides[5]
    add_title(s, "Implementation", "The project is complete but intentionally easy to explain.", None)
    add_card(s, 0.72, 1.65, 3.45, 1.2, "API Layer", "`routes_tests.py` exposes health, tests, runs, and analytics endpoints.", BLUE)
    add_card(s, 4.68, 1.65, 3.45, 1.2, "AI Layer", "`ai_service.py` uses Ollama/Qwen or deterministic fallback rules.", TEAL)
    add_card(s, 8.64, 1.65, 3.45, 1.2, "Execution Layer", "`executor_service.py` runs Playwright actions and captures evidence.", AMBER)
    add_card(s, 2.7, 3.35, 3.45, 1.2, "Data Layer", "`models.py` stores test cases and runs with SQLAlchemy.", RED)
    add_card(s, 6.66, 3.35, 3.45, 1.2, "UI Layer", "A static dashboard is served by FastAPI at `/app`.", BLUE)
    add_footer(s, 6)

    s = slides[6]
    add_title(s, "Result Tracking", "Execution records preserve proof, not just pass/fail labels.", None)
    add_bullets(s, 0.9, 1.8, [
        "TestCase: name, URL, requirement, generated steps, expected result.",
        "TestRun: status, duration, summary, error summary, logs, screenshots.",
        "Analytics: total cases, total runs, pass rate, failures, warnings, recent runs.",
        "Artifacts: screenshots are stored under backend data artifacts.",
    ])
    add_footer(s, 7)

    s = slides[7]
    add_title(s, "Testing", "The backend verification suite confirms the critical demo path.", None)
    add_card(s, 0.9, 1.8, 3.2, 1.25, "Health", "Confirms API availability and model status response.", TEAL)
    add_card(s, 4.8, 1.8, 3.2, 1.25, "Generation", "Creates a saved test case from a URL and requirement.", BLUE)
    add_card(s, 8.7, 1.8, 3.2, 1.25, "Analytics", "Validates counts, run totals, pass rate, and recent run shape.", AMBER)
    add_card(s, 3.45, 4.0, 5.8, 0.85, "Verified", "Current automated test result: 3 passed.", TEAL)
    add_footer(s, 8)

    s = slides[8]
    add_title(s, "Limitations", "The first version is practical, with clear future expansion points.", None)
    add_bullets(s, 0.9, 1.75, [
        "Current version is single-user and local-first.",
        "Locator handling is simple compared with enterprise test frameworks.",
        "AI generation depends on local Ollama availability for best results.",
        "Future scope: PostgreSQL, login system, locator healing, CI/CD, report export, parallel runs, visual regression.",
    ])
    add_footer(s, 9)

    s = slides[9]
    add_title(s, "Conclusion", "The project demonstrates an end-to-end AI-assisted testing workflow.", None)
    add_card(s, 0.95, 1.9, 3.1, 1.3, "Academic Value", "Covers AI, automation, backend APIs, data storage, analytics, and testing.", BLUE)
    add_card(s, 4.85, 1.9, 3.1, 1.3, "Practical Value", "Reduces manual script creation and gives readable execution evidence.", TEAL)
    add_card(s, 8.75, 1.9, 3.1, 1.3, "Demo Value", "Simple workflow: generate, run, inspect logs, view analytics.", AMBER)
    add_card(s, 3.25, 4.2, 6.3, 0.85, "Thank You", "Ready for questions.", INK)
    add_footer(s, 10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
